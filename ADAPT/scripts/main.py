from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
import matlab.engine
import shutil
import os

eng = matlab.engine.start_matlab()
eng.addpath('C:/Users/Tanner Mccoy/OneDrive - Regulus Group/Development/Automation/ADAPT/scripts', nargout=0)

adapt_file = 'raw data/ADAPT  Export on 2020-11-17.xlsx'
eng.ADAPT_Main(adapt_file, nargout=0)

source_dir = 'C:/Users/Tanner Mccoy/OneDrive - Regulus Group/Development/Automation/ADAPT'
target_dir = ''

source_dir = os.getcwd()
target_dir = '/outputs'
file_names = os.listdir(source_dir)

prs = Presentation()
slide = prs.slides.add_slide(slide_layout=6)